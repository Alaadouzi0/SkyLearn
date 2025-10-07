# # ingest.py
# import os
# import json
# import faiss
# import numpy as np
# from sentence_transformers import SentenceTransformer
# import pdfplumber
# from tqdm import tqdm
# from pathlib import Path

# DATA_DIR = Path("../data")
# OUT_DIR = Path("../faiss_index")
# OUT_DIR.mkdir(parents=True, exist_ok=True)
# CHUNKS_FILE = Path("../docs_chunks.jsonl")
# EMB_MODEL_NAME = "all-MiniLM-L6-v2"
# EMB_DIM = 384  # pour all-MiniLM-L6-v2

# def extract_text_from_pdf(pdf_path):
#     texts=[]
#     with pdfplumber.open(pdf_path) as pdf:
#         for page in pdf.pages:
#             txt = page.extract_text()
#             if txt:
#                 texts.append(txt)
#     return "\n".join(texts)

# def chunk_text(text, chunk_size=500, overlap=100):
#     tokens = text.split()
#     chunks=[]
#     i=0
#     while i < len(tokens):
#         chunk = " ".join(tokens[i:i+chunk_size])
#         chunks.append(chunk)
#         i += chunk_size - overlap
#     return chunks

# def main():
#     model = SentenceTransformer(EMB_MODEL_NAME)
#     all_embeddings=[]
#     meta=[]
#     id_counter=0
#     with open(CHUNKS_FILE, "w", encoding="utf-8") as fout:
#         for f in DATA_DIR.iterdir():
#             if f.suffix.lower() in [".pdf", ".txt"]:
#                 print("Processing", f.name)
#                 if f.suffix.lower() == ".pdf":
#                     text = extract_text_from_pdf(str(f))
#                 else:
#                     text = f.read_text(encoding="utf-8")
#                 chunks = chunk_text(text, chunk_size=400, overlap=80)
#                 for idx, chunk in enumerate(chunks):
#                     emb = model.encode(chunk)
#                     all_embeddings.append(emb)
#                     meta_obj = {
#                         "id": id_counter,
#                         "source": f.name,
#                         "chunk_index": idx,
#                         "text": chunk[:4000]
#                     }
#                     fout.write(json.dumps(meta_obj, ensure_ascii=False) + "\n")
#                     meta.append(meta_obj)
#                     id_counter+=1

#     all_embeddings = np.vstack(all_embeddings).astype("float32")
#     index = faiss.IndexFlatIP(EMB_DIM)  # using inner product on normalized vectors
#     # normalize embeddings for cosine similarity:
#     faiss.normalize_L2(all_embeddings)
#     index.add(all_embeddings)
#     faiss.write_index(index, str(OUT_DIR / "faiss.index"))
#     # Save meta list
#     with open(OUT_DIR / "meta.json", "w", encoding="utf-8") as f:
#         json.dump(meta, f, ensure_ascii=False, indent=2)

#     print("Ingestion complete. Indexed:", id_counter)

# if __name__ == "__main__":
#     main()
# ingest.py (version robuste)
import logging
import json
from pathlib import Path
import tempfile

import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from tqdm import tqdm

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger("ingest")

# chemins (adapte si besoin)
DATA_DIR = Path("../data")
OUT_DIR = Path("../faiss_index")
OUT_DIR.mkdir(parents=True, exist_ok=True)
CHUNKS_FILE = Path("../docs_chunks.jsonl")
EMB_MODEL_NAME = "all-MiniLM-L6-v2"

# -----------------------
# Fonctions d'extraction
# -----------------------
def repair_with_pikepdf(path: Path):
    """Essaie de réécrire le PDF avec pikepdf pour corriger certains problèmes."""
    try:
        import pikepdf
    except Exception:
        return None
    try:
        tmp = Path(tempfile.mktemp(suffix=".pdf"))
        with pikepdf.open(str(path)) as pdf:
            pdf.save(str(tmp))
        logger.info("Réparé avec pikepdf: %s -> %s", path.name, tmp.name)
        return tmp
    except Exception as e:
        logger.debug("pikepdf n'a pas pu réparer %s : %s", path.name, e)
        return None

def extract_text_from_pdf(path: Path, ocr_language="fra"):
    """Extraction robuste pour PDF : repair -> pdfplumber -> PyPDF2 -> pdfminer -> OCR."""
    text_parts = []

    # 1) essayer réparation (pikepdf)
    repaired = None
    try:
        repaired = repair_with_pikepdf(path)
        use_path = repaired if repaired else path
    except Exception:
        use_path = path

    # 2) pdfplumber (pdfminer intern)
    try:
        import pdfplumber
        with pdfplumber.open(str(use_path)) as pdf:
            for i, page in enumerate(pdf.pages):
                try:
                    txt = page.extract_text()
                    if txt:
                        text_parts.append(txt)
                except Exception as e:
                    logger.debug("pdfplumber page %d error on %s: %s", i, path.name, e)
        if text_parts:
            return "\n".join(text_parts)
    except Exception as e:
        logger.debug("pdfplumber failed for %s: %s", path.name, e)

    # 3) PyPDF2 fallback
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(use_path))
        for i, page in enumerate(reader.pages):
            try:
                txt = page.extract_text() or ""
                if txt:
                    text_parts.append(txt)
            except Exception as e:
                logger.debug("PyPDF2 page %d error on %s: %s", i, path.name, e)
        if text_parts:
            return "\n".join(text_parts)
    except Exception as e:
        logger.debug("PyPDF2 not available or failed for %s: %s", path.name, e)

    # 4) pdfminer.high_level.extract_text
    try:
        from pdfminer.high_level import extract_text
        txt = extract_text(str(use_path)) or ""
        if txt.strip():
            return txt
    except Exception as e:
        logger.debug("pdfminer.high_level failed for %s: %s", path.name, e)

    # 5) OCR fallback (pdf2image + pytesseract) - optionnel
    try:
        from pdf2image import convert_from_path
        import pytesseract
        images = convert_from_path(str(use_path), dpi=300, fmt="png")
        ocr_texts = []
        for i, img in enumerate(images):
            try:
                o = pytesseract.image_to_string(img, lang=ocr_language)
                if o:
                    ocr_texts.append(o)
            except Exception as e:
                logger.debug("OCR page %d error on %s: %s", i, path.name, e)
        if ocr_texts:
            logger.info("OCR extraction succeeded for %s", path.name)
            return "\n".join(ocr_texts)
    except Exception as e:
        logger.debug("OCR not available or failed for %s: %s", path.name, e)

    # nettoyage fichier temporaire si besoin
    if repaired:
        try:
            repaired.unlink()
        except Exception:
            pass

    # si tout a échoué
    logger.warning("Impossible d'extraire le texte du PDF %s (corrompu ou non supporté)", path.name)
    return ""

def extract_text_from_docx(path: Path):
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        return "\n".join(paragraphs)
    except Exception as e:
        logger.debug("docx extraction failed for %s: %s", path.name, e)
        return ""

def extract_text_from_pptx(path: Path):
    try:
        from pptx import Presentation
        pres = Presentation(str(path))
        texts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = shape.text
                    if t:
                        texts.append(t)
        return "\n".join(texts)
    except Exception as e:
        logger.debug("pptx extraction failed for %s: %s", path.name, e)
        return ""

def extract_text_from_txt(path: Path):
    try:
        return path.read_text(encoding="utf-8")
    except Exception:
        try:
            return path.read_text(encoding="latin-1")
        except Exception as e:
            logger.debug("txt read failed for %s: %s", path.name, e)
            return ""

def extract_text_from_doc_legacy(path: Path):
    """Essayer textract si présent pour .doc (fallback)."""
    try:
        import textract
        text = textract.process(str(path))
        if isinstance(text, bytes):
            try:
                return text.decode("utf-8")
            except Exception:
                return text.decode("latin-1", errors="ignore")
        return str(text)
    except Exception as e:
        logger.debug("textract not available or failed for %s: %s", path.name, e)
        return ""

# -----------------------
# Chunking
# -----------------------
def chunk_text(text, chunk_size=500, overlap=100):
    tokens = text.split()
    if not tokens:
        return []
    chunks = []
    i = 0
    while i < len(tokens):
        chunk = " ".join(tokens[i:i + chunk_size])
        chunks.append(chunk)
        i += chunk_size - overlap
    return chunks

# -----------------------
# Main ingestion
# -----------------------
def main():
    if not DATA_DIR.exists():
        logger.error("DATA_DIR introuvable: %s", DATA_DIR)
        return

    # init modèle
    logger.info("Chargement du modèle d'embeddings: %s", EMB_MODEL_NAME)
    model = SentenceTransformer(EMB_MODEL_NAME)
    emb_dim = model.get_sentence_embedding_dimension()
    logger.info("Embedding dim détectée: %d", emb_dim)

    all_embeddings = []
    meta = []
    id_counter = 0

    # ouvrir fichier chunks en écriture (écrase l'ancien)
    with open(CHUNKS_FILE, "w", encoding="utf-8") as fout:
        files = sorted([p for p in DATA_DIR.rglob("*") if p.is_file()])
        for f in tqdm(files, desc="Fichiers"):
            suffix = f.suffix.lower()
            try:
                logger.info("Processing %s", f.name)
                text = ""
                if suffix == ".pdf":
                    text = extract_text_from_pdf(f)
                elif suffix == ".docx":
                    text = extract_text_from_docx(f)
                elif suffix == ".pptx":
                    text = extract_text_from_pptx(f)
                elif suffix == ".txt":
                    text = extract_text_from_txt(f)
                elif suffix == ".doc":
                    text = extract_text_from_doc_legacy(f)
                else:
                    logger.info("Type non supporté (skipping): %s", f.name)
                    continue

                if not text or not text.strip():
                    logger.warning("Aucun texte extrait pour %s, skip.", f.name)
                    continue

                # chunker et encoder
                chunks = chunk_text(text, chunk_size=400, overlap=80)
                if not chunks:
                    logger.warning("Aucun chunk généré pour %s, skip.", f.name)
                    continue

                # encoder par lot pour efficacité
                embeddings = model.encode(chunks, convert_to_numpy=True, show_progress_bar=False)
                # s'assurer qu'on a bien numpy array 2D
                embeddings = np.array(embeddings, dtype="float32")
                if embeddings.ndim == 1:
                    embeddings = embeddings.reshape(1, -1)

                # vérifier dimension
                if embeddings.shape[1] != emb_dim:
                    logger.warning("Dim embedding inattendue pour %s : %s (attendu %d).", f.name, embeddings.shape, emb_dim)

                # append embeddings + metas
                for idx_chunk, (chunk, emb) in enumerate(zip(chunks, embeddings)):
                    all_embeddings.append(emb)
                    meta_obj = {
                        "id": id_counter,
                        "source": f.name,
                        "chunk_index": idx_chunk,
                        "text": chunk[:4000]
                    }
                    fout.write(json.dumps(meta_obj, ensure_ascii=False) + "\n")
                    meta.append(meta_obj)
                    id_counter += 1

            except Exception as e:
                logger.exception("Erreur sur le fichier %s : %s (skipped)", f.name, e)
                continue

    # si aucun embedding, sortir proprement
    if not all_embeddings:
        logger.error("Aucun embedding généré. Fin du process.")
        return

    # construire numpy array et index FAISS
    arr = np.vstack(all_embeddings).astype("float32")
    # normaliser L2 -> inner product == cosine similarity
    faiss.normalize_L2(arr)
    index = faiss.IndexFlatIP(arr.shape[1])
    index.add(arr)
    faiss.write_index(index, str(OUT_DIR / "faiss.index"))

    # sauver meta
    with open(OUT_DIR / "meta.json", "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

    logger.info("Ingestion complete. Documents indexés: %d embeddings", id_counter)


if __name__ == "__main__":
    main()
